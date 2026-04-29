import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Optional: We could set a dark background but default is fine for auto-generation.
    
    # -------------------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Antigravity"
    
    # Subtitle and body
    subtitle.text = (
        "A Quantum-Ready Behavioral Neuroscience Platform\n"
        "Transforming cognitive intelligence with frontier technologies\n\n"
        "Course: IMBA 25S Frontier Technologies\n"
        "Team: [Your names]\n"
        "Professor: Casimiro Juanes\n"
        "Vision: Predictive Cognitive Intelligence Platform\n"
        "Status: Research phase — seeking guidance on industry direction"
    )
    
    # -------------------------------------------------------------------------
    # SLIDE 2: The Strategic Opportunity
    # -------------------------------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title2 = slide2.shapes.title
    title2.text = "The Problem"

    body2 = slide2.shapes.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "95% of enterprise GenAI pilots fail to deliver measurable ROI (McKinsey / Menlo Ventures)"
    p = tf2.add_paragraph()
    p.text = "AI pilots are fragmented across business units with no unified governance"
    p = tf2.add_paragraph()
    p.text = "EU AI Act (2026-2027) classifies neural data products as high-risk — requiring robust cybersecurity, risk management, and audit trails"
    p = tf2.add_paragraph()
    p.text = "No dominant product connects behavioral data + quantum ML + neuroscience"
    
    p = tf2.add_paragraph()
    p.text = "\nThe Predator Line:"
    p.font.bold = True
    
    p = tf2.add_paragraph()
    p.text = "\"Google is already sitting on the most detailed map of human anxiety ever assembled. The question is who builds the bridge to the clinic before they do.\""
    p.font.italic = True
    
    # -------------------------------------------------------------------------
    # SLIDE 3: The Solution
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    title3 = slide3.shapes.title
    title3.text = "Antigravity Platform"
    
    body3 = slide3.shapes.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "What it is:"
    p = tf3.add_paragraph()
    p.text = "A cybersecurity-first neutral broker platform that aggregates consented behavioral and neural signals, feeds hybrid quantum-ML models, and sells causal cognitive intelligence to governments, pharma, and insurers."
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Core capabilities:"
    
    p = tf3.add_paragraph()
    p.text = "Behavioral data layer — Google Trends + social media emotional signals mapped to population patterns"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Quantum ML layer — IBM Quantum classifies which neural circuits activate (IBM-Inclusive Brains BCI study, June 2025)"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Brain map layer — Neural reference from Meta TRIBE v2 (500h fMRI, 700 subjects, 23% improvement over prior methods)"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Cybersecurity-first — PQC-encrypted, zero-trust, audit-ready, differential privacy"
    p.level = 1

    # -------------------------------------------------------------------------
    # SLIDE 4: Technology & Roadmap
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    title4 = slide4.shapes.title
    title4.text = "Technologies Used + Three Horizons"
    
    body4 = slide4.shapes.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Technologies:"
    
    p = tf4.add_paragraph()
    p.text = "Quantum Computing: IBM Quantum AI (Willow chip, 1000+ qubits) for ML classification"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "AI / GenAI: Hybrid quantum-classical sentiment analysis, cross-linguistic emotional detection"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Cybersecurity: PQC-ready encryption (NIST roadmap), zero-trust architecture, audit logging"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Brain Mapping: Meta TRIBE v2 neural encoding, Harvard/Google brain activity maps"
    p.level = 1
    
    p = tf4.add_paragraph()
    p.text = "Three Horizons:"
    
    p = tf4.add_paragraph()
    p.text = "H1 (Now): Government entry via DSA / EU AI Act compliance tool"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "H2 (12-24 months): Quantum ML biomarker scoring — pharma/insurer contracts"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "H3 (24+ months): New standard of care — pharma trials, cognitive risk pricing, population health"
    p.level = 1

    # -------------------------------------------------------------------------
    # SLIDE 5: Open Decisions / What We Need
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    title5 = slide5.shapes.title
    title5.text = "Decisions We Need From the Professor"
    
    body5 = slide5.shapes.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "We are seeking guidance on three open questions:"
    
    p = tf5.add_paragraph()
    p.text = "Which industry direction should we anchor on?"
    p.level = 1
    
    p = tf5.add_paragraph()
    p.text = "Pharma (Roche, J&J) — trial cohort targeting, biomarker discovery"
    p.level = 2
    p = tf5.add_paragraph()
    p.text = "Insurance (AXA) — cognitive risk scoring, population health pricing"
    p.level = 2
    p = tf5.add_paragraph()
    p.text = "Public Health (NHS, EU Commission) — causal harm dashboard"
    p.level = 2
    
    p = tf5.add_paragraph()
    p.text = "Does the chosen industry count as strategic transformation or disruption per the group brief?"
    p.level = 1
    
    p = tf5.add_paragraph()
    p.text = "How much differentiation is needed between our group presentation and individual reports?"
    p.level = 1
    
    p = tf5.add_paragraph()
    p.text = "Next steps after this meeting:"
    
    p = tf5.add_paragraph()
    p.text = "Lock industry direction and company anchor"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Finalize 6-slide group deck structure & Assign slide presenters"
    p.level = 1

    # Output file path
    out_dir = "/Users/rausan/Library/CloudStorage/OneDrive-Personal/Root/Side Projects/antigravity/group-presentation/vision1/slides"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "vision1_group_pitch_final.pptx")
    
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == '__main__':
    create_presentation()
